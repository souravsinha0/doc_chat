import io
import uuid
from datetime import datetime
from typing import List

from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import pandas as pd
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
import torch

from database import store_document_chunk

# Updated to a better model for semantic retrieval (higher quality than all-MiniLM-L6-v2)
EMBEDDING_MODEL_NAME = "sentence-transformers/all-mpnet-base-v2"
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"Loading embedding model '{EMBEDDING_MODEL_NAME}' on {device}")
embedding_model = SentenceTransformer(EMBEDDING_MODEL_NAME).to(device)

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    length_function=len,
    is_separator_regex=False,
)

def extract_text_from_pdf(file_content: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_content))
    full_text = ""
    for page in reader.pages:
        full_text += page.extract_text() + "\n"
    return full_text

def extract_text_from_docx(file_content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_content))
    full_text = ""
    for para in doc.paragraphs:
        full_text += para.text + "\n"
    return full_text

def extract_text_from_xlsx(file_content: bytes) -> str:
    wb = load_workbook(io.BytesIO(file_content))
    full_text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            full_text += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return full_text

def extract_text_from_csv(file_content: bytes) -> str:
    df = pd.read_csv(io.BytesIO(file_content))
    return df.to_string()

def extract_text_from_pptx(file_content: bytes) -> str:
    prs = Presentation(io.BytesIO(file_content))
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return full_text

# New: Extractor for plain text files (e.g., .py, .txt)
def extract_text_from_txt(file_content: bytes) -> str:
    return file_content.decode('utf-8', errors='ignore')  # Handle any encoding issues

async def process_document(doc_id: uuid.UUID, filename: str, file_content: bytes, file_type: str):
    extractors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'doc': extract_text_from_docx,
        'xlsx': extract_text_from_xlsx,
        'csv': extract_text_from_csv,
        'pptx': extract_text_from_pptx,
        'ppt': extract_text_from_pptx,
        # New: Support for text-based files
        'txt': extract_text_from_txt,
        'py': extract_text_from_txt,
        'md': extract_text_from_txt,  # Optional: Add more extensions as needed
    }
    
    extractor = extractors.get(file_type.lower())
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    full_text = extractor(file_content)
    chunks = text_splitter.split_text(full_text)
    embeddings: List[List[float]] = embedding_model.encode(chunks, convert_to_tensor=False).tolist()
    
    current_time = datetime.utcnow()
    for i, chunk_content in enumerate(chunks):
        await store_document_chunk(
            doc_id=doc_id,
            content=chunk_content,
            embedding=embeddings[i],
            uploaded_at=current_time
        )
    print(f"Processed and stored {len(chunks)} chunks for document {filename} (ID: {doc_id})")

async def process_pdf_document(doc_id: uuid.UUID, filename: str, file_content: bytes):
    await process_document(doc_id, filename, file_content, 'pdf')